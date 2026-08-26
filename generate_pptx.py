import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(15, 23, 42)      # #0f172a
BLUE_ACCENT = RGBColor(59, 130, 246) # #3b82f6
GREEN_ACCENT = RGBColor(16, 185, 129)# #10b981
LIGHT_TEXT = RGBColor(248, 250, 252)# #f8fafc
MUTED_TEXT = RGBColor(148, 163, 184)# #94a3b8
CARD_BG = RGBColor(30, 41, 59)      # #1e293b
CARD_BORDER = RGBColor(51, 65, 85)  # #334155
GOLD_ACCENT = RGBColor(245, 158, 11) # #f59e0b

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, category_text=""):
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    if category_text:
        p0 = tf.paragraphs[0]
        p0.text = category_text.upper()
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = BLUE_ACCENT
        p0.font.name = "Arial"
        
        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.color.rgb = LIGHT_TEXT
        p1.font.name = "Arial"
    else:
        p0 = tf.paragraphs[0]
        p0.text = title_text
        p0.font.size = Pt(26)
        p0.font.bold = True
        p0.font.color.rgb = LIGHT_TEXT
        p0.font.name = "Arial"

# -------------------------------------------------------------
# SLIDE 1: Title Slide
# -------------------------------------------------------------
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide1, DARK_BG)

title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
tf = title_box.text_frame
tf.word_wrap = True

p0 = tf.paragraphs[0]
p0.text = "SATELLITE MANAGEMENT SYSTEM"
p0.font.size = Pt(14)
p0.font.bold = True
p0.font.color.rgb = BLUE_ACCENT

p1 = tf.add_paragraph()
p1.text = "DDD + Hexagonal Architecture + ArchUnit"
p1.font.size = Pt(36)
p1.font.bold = True
p1.font.color.rgb = LIGHT_TEXT

p2 = tf.add_paragraph()
p2.text = "Complete Architectural Blueprint, Codebase Mapping, & Presentation Deck"
p2.font.size = Pt(18)
p2.font.color.rgb = MUTED_TEXT

# -------------------------------------------------------------
# SLIDE 2: DDD — What & Why
# -------------------------------------------------------------
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide2, DARK_BG)
add_header(slide2, "Domain-Driven Design (DDD): What & Why", "Part 1 — Domain-Driven Design")

body_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
tf2 = body_box.text_frame
tf2.word_wrap = True

p = tf2.paragraphs[0]
p.text = "What is Domain-Driven Design (DDD)?"
p.font.bold = True
p.font.size = Pt(20)
p.font.color.rgb = GOLD_ACCENT

items = [
    "A software design approach focused on modeling software to match a complex business domain.",
    "Formulated by Eric Evans, it connects business experts and developers using a Ubiquitous Language."
]
for item in items:
    p = tf2.add_paragraph()
    p.text = "• " + item
    p.font.size = Pt(15)
    p.font.color.rgb = LIGHT_TEXT

p = tf2.add_paragraph()
p.text = "\nWhy Do We Need DDD?"
p.font.bold = True
p.font.size = Pt(20)
p.font.color.rgb = GREEN_ACCENT

items2 = [
    "Eliminates Anemic Domain Models: Data objects are no longer dumb getters/setters; rules live in the domain.",
    "Bridges Communication Gap: Developers and business stakeholders use identical terminology.",
    "Prevents Architectural Decay: Isolates business rules so infrastructure changes don't corrupt domain logic."
]
for item in items2:
    p = tf2.add_paragraph()
    p.text = "• " + item
    p.font.size = Pt(15)
    p.font.color.rgb = LIGHT_TEXT

# -------------------------------------------------------------
# SLIDE 3: DDD Components in Satellite System
# -------------------------------------------------------------
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide3, DARK_BG)
add_header(slide3, "DDD Components in Satellite System", "Part 1 — Domain-Driven Design")

body_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
tf3 = body_box.text_frame
tf3.word_wrap = True

comps = [
    ("Aggregate Root (Satellite.java)", "Core entity managing state machine (REGISTERED -> ACTIVE -> ANOMALY -> DECOMMISSIONED) & invariant enforcement."),
    ("Value Objects (Orbit.java, Telemetry.java)", "Immutable JDK records with self-validation rules (LEO/GEO altitude checks, anomaly battery/temp thresholds)."),
    ("Domain Events (SatelliteLaunchedEvent, AnomalyDetectedEvent)", "Immutable records representing facts that occurred in the domain, recorded internally by the aggregate."),
    ("Strongly-Typed Identity (SatelliteId.java)", "Record wrapping UUID to prevent primitive obsession and guarantee type safety across layers.")
]

for title, desc in comps:
    p = tf3.add_paragraph() if tf3.paragraphs[0].text else tf3.paragraphs[0]
    p.text = "• " + title + ":"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = BLUE_ACCENT
    
    p2 = tf3.add_paragraph()
    p2.text = "   " + desc
    p2.font.size = Pt(14)
    p2.font.color.rgb = LIGHT_TEXT

# -------------------------------------------------------------
# SLIDE 4: DDD Implementation & Challenges
# -------------------------------------------------------------
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide4, DARK_BG)
add_header(slide4, "How to Implement DDD & Challenges", "Part 1 — Domain-Driven Design")

body_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
tf4 = body_box.text_frame
tf4.word_wrap = True

p = tf4.paragraphs[0]
p.text = "How to Implement DDD in Java 21:"
p.font.bold = True
p.font.size = Pt(18)
p.font.color.rgb = GREEN_ACCENT

impls = [
    "Pure Java Domain Model: Zero imports from Spring, JPA, or Jackson in com.satellite.domain.",
    "Java 21 Records: Use records for Value Objects and Domain Events to guarantee immutability.",
    "State Machine Encapsulation: Keep status transition logic private inside Satellite.java."
]
for i in impls:
    p = tf4.add_paragraph()
    p.text = "  1. " + i
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_TEXT

p = tf4.add_paragraph()
p.text = "\nChallenges & Mitigations:"
p.font.bold = True
p.font.size = Pt(18)
p.font.color.rgb = GOLD_ACCENT

challs = [
    "Initial Learning Curve -> Start with clear Bounded Contexts & Ubiquitous Language glossary.",
    "Over-engineering simple CRUD -> Use DDD for core business domains; keep simple lookup tables light.",
    "ORM Entity vs Domain Entity confusion -> Explicitly separate SatelliteJpaEntity from Satellite.java aggregate."
]
for c in challs:
    p = tf4.add_paragraph()
    p.text = "  • " + c
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_TEXT

# -------------------------------------------------------------
# SLIDE 5: DDD Diagram Slide
# -------------------------------------------------------------
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide5, DARK_BG)
add_header(slide5, "DDD Ubiquitous Language & Subdomains Diagram", "Part 1 — Domain-Driven Design")

img_path3 = "/Users/bipin/.gemini/antigravity/scratch/satellite-system/diagram_ddd_ubiquitous_language.png"
if os.path.exists(img_path3):
    slide5.shapes.add_picture(img_path3, Inches(0.8), Inches(1.4), width=Inches(11.7))

# -------------------------------------------------------------
# SLIDE 6: Hexagonal Architecture — What & Why
# -------------------------------------------------------------
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide6, DARK_BG)
add_header(slide6, "Hexagonal Architecture (Ports & Adapters): What & Why", "Part 2 — Hexagonal Architecture")

body_box = slide6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
tf6 = body_box.text_frame
tf6.word_wrap = True

p = tf6.paragraphs[0]
p.text = "What is Hexagonal Architecture?"
p.font.bold = True
p.font.size = Pt(20)
p.font.color.rgb = GOLD_ACCENT

items = [
    "An architectural pattern created by Alistair Cockburn that isolates core application logic inside a Hexagon.",
    "Uses Ports (Interfaces) and Adapters (Implementations) so inputs and outputs can be plugged in seamlessly."
]
for item in items:
    p = tf6.add_paragraph()
    p.text = "• " + item
    p.font.size = Pt(15)
    p.font.color.rgb = LIGHT_TEXT

p = tf6.add_paragraph()
p.text = "\nWhy Do We Need Hexagonal Architecture?"
p.font.bold = True
p.font.size = Pt(20)
p.font.color.rgb = GREEN_ACCENT

items2 = [
    "Decouples Business Logic from External Tools: Database, Web Frameworks, and MQ live outside the Hexagon.",
    "Enables Instant Testing: Domain logic is 100% testable in pure Java without booting Spring or DB (0.01s test execution).",
    "Prevents Technology Lock-in: Swap relational database (H2/Postgres) for MongoDB without changing domain code."
]
for item in items2:
    p = tf6.add_paragraph()
    p.text = "• " + item
    p.font.size = Pt(15)
    p.font.color.rgb = LIGHT_TEXT

# -------------------------------------------------------------
# SLIDE 7: Hexagonal Components in Satellite System
# -------------------------------------------------------------
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide7, DARK_BG)
add_header(slide7, "Hexagonal Components in Satellite System", "Part 2 — Hexagonal Architecture")

body_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
tf7 = body_box.text_frame
tf7.word_wrap = True

hcomps = [
    ("Driving Adapters (Inputs)", "SatelliteController.java (REST API) & TelemetryConsumer.java (Kafka/CLI stream)."),
    ("Driving Ports (Inbound Contracts)", "LaunchSatelliteUseCase.java & UpdateTelemetryUseCase.java (Interfaces defined in domain.port.in)."),
    ("Application Core", "LaunchSatelliteService.java — Orchestrates domain objects via ports in pure Java."),
    ("Driven Ports (Outbound Contracts)", "SatelliteRepository.java & SatelliteEventPublisher.java (Interfaces defined in domain.port.out)."),
    ("Driven Adapters (Outputs)", "SatelliteJpaAdapter.java (SQL), SatelliteMongoAdapter.java (NoSQL), & SatelliteEventPublisherAdapter.java.")
]

for title, desc in hcomps:
    p = tf7.add_paragraph() if tf7.paragraphs[0].text else tf7.paragraphs[0]
    p.text = "• " + title + ":"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = BLUE_ACCENT
    
    p2 = tf7.add_paragraph()
    p2.text = "   " + desc
    p2.font.size = Pt(14)
    p2.font.color.rgb = LIGHT_TEXT

# -------------------------------------------------------------
# SLIDE 8: How to Implement Hexagonal & Challenges
# -------------------------------------------------------------
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide8, DARK_BG)
add_header(slide8, "How to Implement Hexagonal Architecture & Challenges", "Part 2 — Hexagonal Architecture")

body_box = slide8.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
tf8 = body_box.text_frame
tf8.word_wrap = True

p = tf8.paragraphs[0]
p.text = "How to Implement Hexagonal Architecture:"
p.font.bold = True
p.font.size = Pt(18)
p.font.color.rgb = GREEN_ACCENT

himpls = [
    "Package Structure: com.satellite.domain, com.satellite.application, com.satellite.infrastructure.",
    "Inward Dependencies: Infrastructure depends on Application & Domain; Domain depends on NOTHING.",
    "Bean Composition Root: Wire services and adapters in BeanConfig.java inside infrastructure package."
]
for i in himpls:
    p = tf8.add_paragraph()
    p.text = "  1. " + i
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_TEXT

p = tf8.add_paragraph()
p.text = "\nChallenges & Mitigations:"
p.font.bold = True
p.font.size = Pt(18)
p.font.color.rgb = GOLD_ACCENT

hchalls = [
    "Mapping Overhead (Domain <-> DTO <-> JPA Entity) -> Solved with dedicated mappers (SatellitePersistenceMapper).",
    "More interfaces & classes -> Solved by clear package isolation (ports/in, ports/out).",
    "Risk of developers bypassing ports -> Solved by ArchUnit automated architecture build gates!"
]
for c in hchalls:
    p = tf8.add_paragraph()
    p.text = "  • " + c
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_TEXT

# -------------------------------------------------------------
# SLIDE 9: Herberto Graça Explicit Architecture Diagram Slide
# -------------------------------------------------------------
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide9, DARK_BG)
add_header(slide9, "Herberto Graça Explicit Architecture Mapping", "Part 2 — Hexagonal Architecture")

img_path0 = "/Users/bipin/.gemini/antigravity/scratch/satellite-system/diagram_explicit_architecture.png"
if os.path.exists(img_path0):
    slide9.shapes.add_picture(img_path0, Inches(0.8), Inches(1.4), width=Inches(11.7))

# -------------------------------------------------------------
# SLIDE 10: Multi-Database Adapter Swapping Demo Diagram Slide
# -------------------------------------------------------------
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide10, DARK_BG)
add_header(slide10, "Multi-Database Adapter Swapping (SQL vs NoSQL)", "Part 2 — Hexagonal Architecture")

img_path4 = "/Users/bipin/.gemini/antigravity/scratch/satellite-system/diagram_adapter_swapping.png"
if os.path.exists(img_path4):
    slide10.shapes.add_picture(img_path4, Inches(0.8), Inches(1.4), width=Inches(11.7))

# -------------------------------------------------------------
# SLIDE 11: ArchUnit — What & Why
# -------------------------------------------------------------
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide11, DARK_BG)
add_header(slide11, "ArchUnit Architecture Enforcement: What & Why", "Part 3 — ArchUnit Enforcement")

body_box = slide11.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
tf11 = body_box.text_frame
tf11.word_wrap = True

p = tf11.paragraphs[0]
p.text = "What is ArchUnit?"
p.font.bold = True
p.font.size = Pt(20)
p.font.color.rgb = GOLD_ACCENT

items = [
    "A Java architecture testing library that checks Java bytecode (.class files) using fluent JUnit 5 tests.",
    "Analyzes package dependencies, class naming conventions, annotation rules, and layer isolation."
]
for item in items:
    p = tf11.add_paragraph()
    p.text = "• " + item
    p.font.size = Pt(15)
    p.font.color.rgb = LIGHT_TEXT

p = tf11.add_paragraph()
p.text = "\nWhy Do We Need ArchUnit?"
p.font.bold = True
p.font.size = Pt(20)
p.font.color.rgb = GREEN_ACCENT

items2 = [
    "Automates Code Reviews: Architecture rules are checked on every commit in CI/CD pipelines.",
    "Prevents Architectural Erosion: Prevents quick shortcuts (e.g. Controller calling DB Repository directly).",
    "Runs as a Pre-Packaging Build Gate: Executes during mvn test — FAILS THE BUILD before packaging or deployment!"
]
for item in items2:
    p = tf11.add_paragraph()
    p.text = "• " + item
    p.font.size = Pt(15)
    p.font.color.rgb = LIGHT_TEXT

# -------------------------------------------------------------
# SLIDE 12: ArchUnit Rules in Satellite System
# -------------------------------------------------------------
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide12, DARK_BG)
add_header(slide12, "ArchUnit Rules & Build Gates in Satellite System", "Part 3 — ArchUnit Enforcement")

body_box = slide12.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
tf12 = body_box.text_frame
tf12.word_wrap = True

arules = [
    ("Domain Purity Rule", "noClasses().that().resideInAPackage('..domain..').should().dependOnClassesThat().resideInAPackage('..infrastructure..')"),
    ("No Spring/JPA in Domain", "noClasses().that().resideInAPackage('..domain.model..').should().beAnnotatedWith('..Entity')"),
    ("Controller Port Constraint", "noClasses().that().resideInAPackage('..rest..').should().dependOnClassesThat().resideInAPackage('..application..')"),
    ("Layered Architecture Rule", "layeredArchitecture().layer('Domain').layer('Application').layer('Infrastructure')")
]

for title, code in arules:
    p = tf12.add_paragraph() if tf12.paragraphs[0].text else tf12.paragraphs[0]
    p.text = "• " + title + ":"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = BLUE_ACCENT
    
    p2 = tf12.add_paragraph()
    p2.text = "   " + code
    p2.font.size = Pt(13)
    p2.font.color.rgb = LIGHT_TEXT

# -------------------------------------------------------------
# SLIDE 13: How to Implement ArchUnit & Challenges
# -------------------------------------------------------------
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide13, DARK_BG)
add_header(slide13, "How to Implement ArchUnit & Challenges", "Part 3 — ArchUnit Enforcement")

body_box = slide13.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
tf13 = body_box.text_frame
tf13.word_wrap = True

p = tf13.paragraphs[0]
p.text = "How to Implement ArchUnit:"
p.font.bold = True
p.font.size = Pt(18)
p.font.color.rgb = GREEN_ACCENT

aimpls = [
    "Add Maven Dependency: archunit-junit5 (v1.3.0) in pom.xml.",
    "Write Test Class: Annotate with @AnalyzeClasses(packages = 'com.satellite').",
    "Define Rules: Write static @ArchTest rules using ArchRuleDefinition DSL.",
    "Lifecycle Execution: Executes automatically post-compilation during mvn test before packaging!"
]
for i in aimpls:
    p = tf13.add_paragraph()
    p.text = "  1. " + i
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_TEXT

p = tf13.add_paragraph()
p.text = "\nChallenges & Mitigations:"
p.font.bold = True
p.font.size = Pt(18)
p.font.color.rgb = GOLD_ACCENT

achalls = [
    "Retrofitting onto legacy codebases -> Solved with FreezingArchRule to freeze existing violations.",
    "Refactoring package structures -> Update package string variables in HexagonalArchitectureTest.java."
]
for c in achalls:
    p = tf13.add_paragraph()
    p.text = "  • " + c
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_TEXT

# -------------------------------------------------------------
# SLIDE 14: Execution & Live Demo Steps
# -------------------------------------------------------------
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide14, DARK_BG)
add_header(slide14, "Execution & Live Demo Steps", "Summary & Live Demo")

body_box = slide14.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
tf14 = body_box.text_frame
tf14.word_wrap = True

p = tf14.paragraphs[0]
p.text = "1. Run Automated Test Suite & ArchUnit Build Gates:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = GOLD_ACCENT

p = tf14.add_paragraph()
p.text = "  $ mvn clean test   # Executes 37 unit, integration & ArchUnit tests (0 failures)"
p.font.size = Pt(14)
p.font.color.rgb = LIGHT_TEXT

p = tf14.add_paragraph()
p.text = "\n2. Run Spring Boot Application & API Verification:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = GREEN_ACCENT

p = tf14.add_paragraph()
p.text = "  $ mvn spring-boot:run\n  $ curl -X POST http://localhost:8080/api/v1/satellites -d '{...}'"
p.font.size = Pt(14)
p.font.color.rgb = LIGHT_TEXT

p = tf14.add_paragraph()
p.text = "\n3. Live Layer Conflict Demo (ArchUnit Build Failure):"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = BLUE_ACCENT

p = tf14.add_paragraph()
p.text = "  • Inject LaunchSatelliteService into SatelliteController.java\n  • Run mvn test -> ArchUnit FAILS THE BUILD with explicit violation error report!"
p.font.size = Pt(14)
p.font.color.rgb = LIGHT_TEXT

prs.save("/Users/bipin/.gemini/antigravity/scratch/satellite-system/Satellite_Architecture_DDD_Hexagonal_ArchUnit_FINAL.pptx")

print("Successfully generated FINAL PowerPoint presentation deck: Satellite_Architecture_DDD_Hexagonal_ArchUnit_FINAL.pptx")
